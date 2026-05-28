#!/usr/bin/env python3
"""Build Lineup.pptx — a 3-slide pitch deck. Run: python3 scripts/build_pptx.py"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
QR = os.path.join(HERE, "repo-qr.png")
OUT = os.path.join(HERE, "Lineup.pptx")
REPO = "https://github.com/kaiser-data/lineup"

INK   = RGBColor(0x1A, 0x1A, 0x1A)
MUTE  = RGBColor(0x6B, 0x72, 0x80)
RED   = RGBColor(0xEF, 0x44, 0x44)
BLUE  = RGBColor(0x25, 0x63, 0xEB)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
BG    = RGBColor(0xFA, 0xFA, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, l, t, w, h, fill=None, line=None, line_w=1.0, round_=True):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        l, t, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6):
    """runs: list of paragraphs; each paragraph is list of (txt, size, bold, color)."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        for (txt, size, bold, color) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = "Helvetica Neue"
    return tb


# ───────────────────────── Slide 1 ─────────────────────────
s = prs.slides.add_slide(BLANK); bg(s)
box(s, Inches(0), Inches(0), Inches(0.28), SH, fill=RED, round_=False)
text(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(2.0), [
    [("Lineup", 60, True, INK)],
    [("Not a Luma replacement — the fast lane next to it.", 24, False, MUTE)],
])
text(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.5), [
    [("Keep Luma for the real event page. But the moment you need to ", 20, False, INK),
     ("announce it", 20, True, INK),
     (" — an email, a Slack post, an Instagram story, a “who’s speaking” "
      "card — you’re back to juggling Canva + a QR generator + a calendar file.", 20, False, INK)],
])
box(s, Inches(0.9), Inches(4.3), Inches(11.5), Inches(0.95), fill=WHITE, line=RED, line_w=1.5)
text(s, Inches(1.2), Inches(4.45), Inches(11.0), Inches(0.7), [
    [("“Announce Friday’s panel — speakers Lea, Marco and Priya — accent red.”",
      22, True, RED)],
], anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(1.5), [
    [("→ a ready-to-post ", 20, False, INK), ("event card", 20, True, INK),
     (" · an ", 20, False, INK), ("identity badge per speaker", 20, True, INK),
     (" (avatar + role + contact QR) · an ", 20, False, INK),
     ("RSVP QR", 20, True, INK), (" · a calendar ", 20, False, INK),
     (".ics", 20, True, INK), (". Drop them into the email or the post. ", 20, False, INK),
     ("~10 seconds.", 20, True, RED)],
])

# ───────────────────────── Slide 2 ─────────────────────────
s = prs.slides.add_slide(BLANK); bg(s)
box(s, Inches(0), Inches(0), Inches(0.28), SH, fill=BLUE, round_=False)
text(s, Inches(0.9), Inches(0.55), Inches(11.5), Inches(1.2), [
    [("One primitive, many use cases", 40, True, INK)],
    [("people + a moment + a scannable code", 20, False, MUTE)],
])

rows = [
    ("\U0001F4E3", "Announce an event", "Card + QR + .ics to paste into an email or social post"),
    ("\U0001F3A4", "Share who’s speaking", "A badge per speaker — avatar, role, contact QR"),
    ("\U0001F5F3️", "Voting / polls", "A unique code sent to each person — one ballot QR per voter"),
    ("\U0001FAAA", "Conference ID / access", "Printable badge, role, check-in QR"),
    ("\U0001F91D", "Networking", "Badge QR drops your contact straight into a phone"),
    ("\U0001F3C6", "Team rosters / crews", "A collectible “trading card” per member"),
]
top = Inches(1.95); rh = Inches(0.78); gap = Inches(0.08)
for i, (emoji, title, desc) in enumerate(rows):
    y = Emu(int(top) + i * int(rh + gap))
    box(s, Inches(0.9), y, Inches(11.5), rh, fill=WHITE, line=RGBColor(0xE4,0xE4,0xE7), line_w=1.0)
    text(s, Inches(1.1), y, Inches(0.8), rh, [[(emoji, 22, False, INK)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(1.9), y, Inches(3.6), rh, [[(title, 19, True, INK)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(5.6), y, Inches(6.6), rh, [[(desc, 16, False, MUTE)]], anchor=MSO_ANCHOR.MIDDLE)

text(s, Inches(0.9), Inches(6.95), Inches(11.5), Inches(0.5), [
    [("Complements Luma, Mailchimp, your CMS — the 10-second asset generator inside the chat you’re already in.",
      15, True, BLUE)],
], align=PP_ALIGN.CENTER)

# ───────────────────────── Slide 3 ─────────────────────────
s = prs.slides.add_slide(BLANK); bg(s)
box(s, Inches(0), Inches(0), Inches(0.28), SH, fill=GREEN, round_=False)
text(s, Inches(0.9), Inches(0.55), Inches(8.0), Inches(1.0), [
    [("How & where it runs", 40, True, INK)],
])
text(s, Inches(0.9), Inches(1.55), Inches(7.6), Inches(2.2), [
    [("The stack is the moat — each piece does one thing well:", 19, True, INK)],
    [("Skybridge (Alpic)", 18, True, INK), (" — a real React UI inside ChatGPT + tunnel + deploy", 18, False, INK)],
    [("qrcode · DiceBear · ics", 18, True, INK), (" — tiny deterministic generators", 18, False, INK)],
    [("The LLM", 18, True, INK), (" — parses your sentence into fields, writes the copy", 18, False, INK)],
])
# mono flow box
box(s, Inches(0.9), Inches(3.85), Inches(7.6), Inches(1.7), fill=RGBColor(0x1E,0x29,0x32), line=None)
tb = s.shapes.add_textbox(Inches(1.1), Inches(4.0), Inches(7.2), Inches(1.4))
tf = tb.text_frame; tf.word_wrap = True
flow = [
    "ChatGPT  --/mcp (No Auth)-->  Alpic tunnel",
    "                -->  Skybridge server (localhost:3000)",
    "                      |- generate-lineup tool",
    "                      |- qrcode . DiceBear . ics",
    "                      `- React view (in ChatGPT)",
]
for i, ln in enumerate(flow):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = ln
    r.font.size = Pt(12); r.font.name = "Courier New"
    r.font.color.rgb = RGBColor(0x8C,0xE9,0xB0)
text(s, Inches(0.9), Inches(5.75), Inches(7.6), Inches(1.5), [
    [("Now:", 16, True, INK), (" runs on the laptop via an Alpic tunnel for the live demo.", 16, False, INK)],
    [("Submission:", 16, True, INK), (" npm run deploy → Alpic Cloud (free). Pure-JS / edge-safe → Cloudflare Workers fallback. No accounts, no scraping.", 16, False, INK)],
])

# QR + link panel (right)
box(s, Inches(8.9), Inches(1.55), Inches(3.5), Inches(5.0), fill=WHITE, line=RGBColor(0xE4,0xE4,0xE7), line_w=1.5)
text(s, Inches(8.9), Inches(1.75), Inches(3.5), Inches(0.5), [
    [("Get the code", 18, True, INK)],
], align=PP_ALIGN.CENTER)
if os.path.exists(QR):
    s.shapes.add_picture(QR, Inches(9.35), Inches(2.35), height=Inches(2.6))
text(s, Inches(8.9), Inches(5.1), Inches(3.5), Inches(0.6), [
    [("github.com/", 15, False, MUTE)],
    [("kaiser-data/lineup", 16, True, GREEN)],
], align=PP_ALIGN.CENTER)
text(s, Inches(8.9), Inches(6.0), Inches(3.5), Inches(0.5), [
    [("Scan to clone · demo it live", 13, False, MUTE)],
], align=PP_ALIGN.CENTER)

text(s, Inches(0.9), Inches(7.0), Inches(11.5), Inches(0.4), [
    [("Lineup: the fastest way to turn a sentence into something everyone can hold in their hand.",
      16, True, INK)],
], align=PP_ALIGN.CENTER)

prs.save(OUT)
print("wrote", OUT)
